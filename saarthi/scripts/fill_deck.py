"""
Fill the IDBI Innovate prototype submission deck for SAARTHI — Problem Statement 4.
Adds concise text in the existing style, inserts the process-flow / architecture
diagrams and the prototype screenshot, and fills the links slide. Optional pages
(Wireframes, Cost) are left untouched. Format and existing images are preserved.
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

ROOT = "/home/Debz/Hackathon/IDBI_Hackathon"
SRC = f"{ROOT}/Prototype Submission Deck _ IDBI Innovate.pptx"
FIGS = f"{ROOT}/deck_figs"
SHOT = f"{ROOT}/Saarthi.png"
INK = RGBColor(0x1F, 0x29, 0x37)
LINKBLUE = RGBColor(0x1D, 0x4E, 0xD8)

prs = Presentation(SRC)


def add_box(slide, top, paragraphs, size=12, left=0.4, width=9.2, space_after=6):
    """paragraphs: list of paragraphs; each is a list of (text, bold, color?) runs."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.4 - top))
    tf = tb.text_frame; tf.word_wrap = True
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.line_spacing = 1.05
        for run in runs:
            text, bold = run[0], run[1]
            col = run[2] if len(run) > 2 else INK
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col
    return tb


def fit(img, max_w=9.2, max_h=3.7):
    w, h = Image.open(img).size
    ar = w / h
    W, H = max_w, max_w / ar
    if H > max_h:
        H, W = max_h, max_h * ar
    return W, H, (10 - W) / 2, 1.6  # width, height, centered-left, top


def add_pic(slide, img):
    W, H, L, T = fit(img)
    slide.shapes.add_picture(img, Inches(L), Inches(T), Inches(W), Inches(H))


def append_after_labels(slide, mapping, size=12):
    """For a text box whose paragraphs start with known labels, append a value run."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            txt = "".join(r.text for r in p.runs).strip()
            for label, (value, link) in mapping.items():
                if txt.lower().startswith(label.lower()):
                    r = p.add_run(); r.text = value
                    r.font.size = Pt(size); r.font.bold = False
                    r.font.color.rgb = LINKBLUE if link else INK
                    if link:
                        r.hyperlink.address = value.strip()


S = prs.slides

# ── Slide 0 — Team details ────────────────────────────────────────────────
append_after_labels(S[0], {
    "Team name:": ("  SAARTHI", None),
    "Team leader name:": ("  Koushik Deb", None),
    "Problem Statement:": ("  Problem Statement 4 – Default Prediction Model", None),
}, size=15)

# ── Slide 1 — Brief about the idea ────────────────────────────────────────
add_box(S[1], 1.55, [
    [("SAARTHI is an MSME loan default early-warning system. A lender uploads a loan "
      "dataset (CSV); SAARTHI auto-maps the columns, trains a calibrated model, and for "
      "every loan predicts the probability of default over the next 12 months.", False)],
    [("But it never stops at a number. Each loan comes with a plain-English reason, a "
      "recommended action to reduce the risk, and a fairness check — and every explanation "
      "is verified by a second AI, so it cannot make things up.", False)],
    [("In short: it turns a raw risk score into a decision a credit officer can act on.", True)],
], size=14, space_after=10)

# ── Slide 2 — Opportunities ───────────────────────────────────────────────
add_box(S[2], 2.05, [
    [("How it is different:  ", True),
     ("Most tools give only a risk score. SAARTHI adds a verified plain-English reason, a "
      "12-month timeline, and a concrete fix for every loan — all in one common vocabulary, "
      "so any two loans are comparable.", False)],
    [("How it solves the problem:  ", True),
     ("The ML model gives a calibrated default probability; SHAP finds the drivers; an AI "
      "writes the reason and the action; a second 'faithfulness judge' checks it against the "
      "evidence; and a fairness audit guards against bias.", False)],
    [("USP:  ", True),
     ("Trustworthy, explained output — never a bare number. The AI can describe the model but "
      "is not allowed to invent it, and the single move that lowers the risk is shown as "
      "before to after.", False)],
], size=12.5, space_after=10)

# ── Slide 3 — List of features ────────────────────────────────────────────
feats = [
    "Auto column mapping — any CSV to a fixed schema, with user override",
    "Calibrated 12-month probability of default (PD) for every loan",
    "Early-warning risk curve that flags the onset month",
    "SHAP-based reasons in a fixed, comparable vocabulary (common interpretation framework)",
    "Plain-English explanation, verified by a faithfulness judge (anti-hallucination)",
    "Recommended action with before to after PD (recourse)",
    "Difference-aware fairness audit — protected attributes are never used to predict",
    "Portfolio dashboard + per-loan drill-down + model trace",
    "Multi-provider AI gateway with automatic fail-over; honest labelling of estimates",
]
add_box(S[3], 1.5, [[("•  " + f, False)] for f in feats], size=12, space_after=5)

# ── Slide 4 — Process flow ────────────────────────────────────────────────
add_pic(S[4], f"{FIGS}/process_flow.png")

# ── Slide 6 — Architecture ────────────────────────────────────────────────
add_pic(S[6], f"{FIGS}/architecture.png")

# ── Slide 7 — Technologies ────────────────────────────────────────────────
tech = [
    ("Frontend:  ", "React, Vite, TypeScript, Tailwind CSS, Recharts"),
    ("Backend:  ", "Python, Flask, Gunicorn (async job manager)"),
    ("Machine Learning:  ", "LightGBM, scikit-learn, SHAP, lifelines (survival), fairlearn, pandas, polars"),
    ("AI / LLM:  ", "DeepSeek, Mistral, OpenRouter, Google Gemini (OpenAI-compatible) · pydantic + json-repair"),
    ("Hosting:  ", "Caddy (automatic HTTPS), systemd, DuckDNS, Ubuntu server"),
]
add_box(S[7], 1.55, [[(a, True), (b, False)] for a, b in tech], size=13, space_after=11)

# ── Slide 9 — Snapshots of the prototype ──────────────────────────────────
add_pic(S[9], SHOT)

# ── Slide 10 — Performance / benchmarking ─────────────────────────────────
perf = [
    "Accuracy (AUC-ROC): 0.97 on the real SBA loan book (~900K loans); 0.68–0.79 on German & MSME sets",
    "Calibration (ECE): about 0.05 — the probabilities are reliable, not just a ranking",
    "Explanation faithfulness: ~90–100% verified by the judge across test loans",
    "Scale & speed: scores 148K+ loans; trains in seconds; a per-loan explanation takes ~10–20s",
    "Fairness: demographic-parity & equalized-odds gaps reported per protected attribute — Pass on test data",
    "Reliability: 33/33 automated end-to-end checks and 12/12 live API checks pass",
]
add_box(S[10], 1.5, [[("•  " + f, False)] for f in perf], size=12, space_after=6)

# ── Slide 11 — Additional details / future development ────────────────────
future = [
    "Integration with bank core / loan-origination systems and scheduled portfolio re-scoring",
    "Role-based access and on-premise / private deployment",
    "More models (TabPFN, XGBoost) and automatic data-drift monitoring",
    "Regulator-ready audit logs and exportable reason reports",
    "Research paper on the capability-diverse faithfulness judge",
]
add_box(S[11], 1.5, [[("•  " + f, False)] for f in future], size=12.5, space_after=7)

# ── Slide 12 — Links ──────────────────────────────────────────────────────
append_after_labels(S[12], {
    "GitHub Public Repository": (":  https://github.com/DevDaring/IDBI_SAARTHI", True),
    "Demo Video Link": (":  https://youtu.be/-HglTjzO5qc", True),
    "Final Product Link": (":  https://koushikdeb.duckdns.org", True),
}, size=12)

# ── save (back up original first) ─────────────────────────────────────────
shutil.copy(SRC, f"{ROOT}/Prototype Submission Deck _ IDBI Innovate (original backup).pptx")
prs.save(SRC)
print("Deck filled and saved:", SRC)
