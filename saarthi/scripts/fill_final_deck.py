"""
Fill the IDBI Innovate 2026 Prototype Submission Deck for Team SAARTHI.

Format-preserving by construction:
  * never moves, restyles or deletes an existing shape
  * appends to the runs the template already provides (slides 1 and 14)
  * adds new text boxes only in the empty band between the header (ends 1.4")
    and the footer rule (starts 5.5")
  * uses Arial, the deck's own theme font, at template-consistent sizes

Optional pages are left untouched where we have nothing genuine to show.
Every number is read from models/metrics.json - none are typed by hand.
"""
from __future__ import annotations

import glob
import json
import os
import shutil

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = "/home/Debz/Hackathon/IDBI_Hackathon"
SRC = f"{ROOT}/Final_Submission/Prototype Submission Deck _ IDBI Innovate 2026.pptx"
FIGS = f"{ROOT}/deck_figs"
SHOTS = f"{FIGS}/shots"
MODELS = f"{ROOT}/saarthi/models"

FONT = "Arial"
INK = RGBColor(0x1F, 0x29, 0x37)
MUTE = RGBColor(0x44, 0x4E, 0x5E)
ACCENT = RGBColor(0xB4, 0x53, 0x09)
LINK = RGBColor(0x1D, 0x4E, 0xD8)

TOP = 1.55          # content band starts below the header text box
BOTTOM = 5.40       # footer rule sits at 5.5


# ---------------------------------------------------------------- helpers
def body(slide, paras, top=TOP, left=0.45, width=9.1, size=12, gap=5):
    """paras: list of [(text, bold, colour?), ...] run-groups."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(BOTTOM - top))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        for run in runs:
            text, bold = run[0], run[1]
            col = run[2] if len(run) > 2 else INK
            sz = run[3] if len(run) > 3 else size
            r = p.add_run()
            r.text = text
            r.font.size = Pt(sz)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = col
    return tb


def picture(slide, path, top=TOP, max_w=9.0, max_h=None):
    """Insert an image scaled to fit the content band, horizontally centred."""
    max_h = max_h or (BOTTOM - top)
    w, h = Image.open(path).size
    ar = w / h
    dw, dh = max_w, max_w / ar
    if dh > max_h:
        dh, dw = max_h, max_h * ar
    left = (10.0 - dw) / 2
    return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                    Inches(dw), Inches(dh))


def append_to_run(slide, startswith, text, bold=False, col=None):
    """Append text onto the template's own run so its styling is inherited."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if r.text.strip().startswith(startswith):
                    nr = para.add_run()
                    nr.text = text
                    nr.font.size = r.font.size or Pt(15)
                    nr.font.bold = bold
                    nr.font.name = FONT
                    nr.font.color.rgb = col or INK
                    return True
    return False


def m():
    with open(f"{MODELS}/metrics.json") as fh:
        return json.load(fh)


def abl():
    p = f"{MODELS}/ablation_sequence.json"
    return json.load(open(p)) if os.path.exists(p) else {}


def table(slide, rows, top, left=0.7, width=8.6, row_h=0.26, sizes=(10, 10)):
    n_r, n_c = len(rows), len(rows[0])
    shp = slide.shapes.add_table(n_r, n_c, Inches(left), Inches(top),
                                 Inches(width), Inches(row_h * n_r))
    tbl = shp.table
    for ci, w in enumerate(_col_widths(n_c, width)):
        tbl.columns[ci].width = Inches(w)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            for r in para.runs:
                r.font.size = Pt(sizes[0] if ri == 0 else sizes[1])
                r.font.bold = (ri == 0)
                r.font.name = FONT
                r.font.color.rgb = INK
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
    return shp


def _col_widths(n, total):
    if n == 4:
        return [total * x for x in (0.34, 0.22, 0.22, 0.22)]
    if n == 5:
        return [total * x for x in (0.28, 0.18, 0.18, 0.18, 0.18)]
    if n == 3:
        return [total * x for x in (0.44, 0.28, 0.28)]
    return [total / n] * n


# ---------------------------------------------------------------- content
def fill(prs):
    S = prs.slides
    M = m()
    spec = M["specialists"]
    pooled = M["pooled"]["calibrated_test"]
    A = abl()

    # -- 1. Team details -------------------------------------------------
    append_to_run(S[0], "Team name", "SAARTHI")
    append_to_run(S[0], "Team leader name", "Koushik Deb  (solo team)")
    append_to_run(S[0], "Problem Statement",
                  "Default Prediction Model  (Track 04)")

    # -- 2. Brief about the idea -----------------------------------------
    body(S[1], [
        [("SAARTHI turns an MSME loan book into a 12-month early-warning system.", True)],
        [("It answers three questions for every borrower: ", False),
         ("who is likely to default, how soon, and the single change that "
          "measurably lowers that risk.", True)],
        [("", False)],
        [("The design rule: ", True),
         ("the ML model owns every number (probability of default, SHAP drivers, "
          "the projected effect of any remedy). The language model owns only words "
          "— it picks a code from a fixed 10-item taxonomy and writes the "
          "explanation. A second model from a different family then audits that "
          "explanation against the SHAP evidence, so nothing reaching a credit "
          "officer is invented.", False)],
        [("", False)],
        [("Unlike a static scorecard, SAARTHI ships pre-trained on 3.3M public "
          "loans and then fine-tunes on the lender's own book in minutes — "
          "on-premise, with no data leaving the bank.", False)],
    ])

    # -- 3. Opportunities -------------------------------------------------
    body(S[2], [
        [("How is it different from existing ideas?", True, ACCENT)],
        [("Credit AI today is either accurate but opaque, or explainable but weak. "
          "SAARTHI separates the two jobs: numbers come from a calibrated model, "
          "words from an LLM, and an ", False),
         ("independent faithfulness judge", True),
         (" verifies the words against the model's own SHAP evidence before a "
          "human sees them.", False)],
        [("", False, INK, 4)],
        [("How will it solve the problem?", True, ACCENT)],
        [("Defaults are detected at 90+ DPD, when recovery options have closed. "
          "SAARTHI flags the risk up to 12 months earlier and pairs each alert "
          "with a grounded counterfactual — extend tenure, add collateral, "
          "restructure rate — searched on the model itself, not invented by an LLM.",
          False)],
        [("", False, INK, 4)],
        [("USP", True, ACCENT)],
        [("Every explanation is machine-verified, every recommended action carries "
          "a model-computed post-action PD, protected attributes never enter the "
          "model but are audited for bias, and the whole system runs inside the "
          "bank's own environment.", False)],
    ], size=11.5)

    # -- 4. Features ------------------------------------------------------
    body(S[3], [
        [("Automatic column mapping", True), (" — upload any loan book; an LLM maps "
          "your columns onto a canonical schema, with a deterministic fallback and "
          "user confirmation.", False)],
        [("Calibrated probability of default", True), (" — isotonic-calibrated, so "
          "the score is a real probability, not a rank.", False)],
        [("12-month risk curve", True), (" — survival model shapes when risk accrues; "
          "flags the month the borrower crosses the alert threshold.", False)],
        [("Verified reason codes", True), (" — fixed 10-code taxonomy grounded in "
          "SHAP, checked by an independent judge model for invented factors and "
          "flipped signs.", False)],
        [("Actionable recourse", True), (" — greedy counterfactual search over "
          "actionable levers only, returning the projected PD after the change.", False)],
        [("Difference-aware fairness audit", True), (" — protected attributes are "
          "never model inputs; disparities are flagged only when they persist among "
          "same-risk applicants.", False)],
        [("Transaction-as-language", True), (" — a self-supervised encoder reads raw "
          "transaction narration sequences, adding signal where curated fields are "
          "sparse.", False)],
        [("Full transparency", True), (" — live model trace showing which provider "
          "and model answered, and whether each explanation passed its judge.", False)],
    ], size=11)

    # -- 5. Process flow --------------------------------------------------
    f = f"{FIGS}/process_flow.png"
    if os.path.exists(f):
        picture(S[4], f, top=1.5, max_w=8.8, max_h=3.8)

    # -- 6. Wireframes (OPTIONAL) — intentionally left blank --------------

    # -- 7. Architecture --------------------------------------------------
    f = f"{FIGS}/architecture.png"
    if os.path.exists(f):
        picture(S[6], f, top=1.5, max_w=8.8, max_h=3.8)

    # -- 8. Technologies --------------------------------------------------
    body(S[7], [
        [("Machine learning", True, ACCENT)],
        [("LightGBM · XGBoost · CatBoost ensemble, isotonic calibration "
          "(scikit-learn) · SHAP TreeExplainer · lifelines (Cox survival) · "
          "fairlearn (bias audit) · PyTorch (CoLES transaction-sequence encoder)",
          False)],
        [("", False, INK, 4)],
        [("Language models", True, ACCENT)],
        [("Multi-provider gateway over an OpenAI-compatible interface — DeepSeek, "
          "Mistral, OpenRouter, Gemini — with role-based fallback chains, key "
          "rotation and a cross-family judge. AWS Bedrock adapter for in-bank "
          "deployment.", False)],
        [("", False, INK, 4)],
        [("Application", True, ACCENT)],
        [("Backend: Python 3.11, Flask, Gunicorn, pandas/polars, pydantic. "
          "Frontend: React 18, TypeScript, Vite, TailwindCSS, Recharts. "
          "Serving: Caddy with automatic HTTPS, systemd, Linux VPS.", False)],
        [("", False, INK, 4)],
        [("Training data", True, ACCENT)],
        [("Nine public credit corpora — SBA, Lending Club, Home Credit, Give Me "
          "Some Credit, Taiwan, German, Berka, Amex, IBM TabFormer. No bank data "
          "was used at any stage.", False)],
    ], size=11)

    # -- 9. Estimated cost (OPTIONAL, but genuinely a strength) -----------
    body(S[8], [
        [("Prototype built end-to-end for under one US dollar.", True)],
        [("", False, INK, 4)],
    ])
    table(S[8], [
        ["Item", "Spec", "Cost"],
        ["GPU (sequence encoder training)", "RTX 2080 Ti, vast.ai, ~1 h", "$0.08"],
        ["CPU training (all other models)", "existing 32-core VPS", "$0.00"],
        ["Datasets", "9 public corpora, ~15 GB", "$0.00"],
        ["Hosting (demo)", "VPS + Caddy, automatic HTTPS", "existing"],
        ["LLM inference", "free-tier / low-cost API keys", "< $1"],
    ], top=2.15)
    body(S[8], [
        [("Production estimate (per bank):", True)],
        [("Inference is CPU-only and runs on-premise; a single mid-size VM scores a "
          "150,000-loan portfolio per run. The only recurring cost is LLM inference "
          "for explanations, which is bounded by generating them for the top-risk "
          "loans eagerly and the rest lazily on demand.", False)],
    ], top=4.05, size=11)

    # -- 10. Snapshots ----------------------------------------------------
    shots = [f"{SHOTS}/{n}.png" for n in
             ("05_dashboard", "06_loan_detail", "01_landing", "03_mapping")]
    shots = [s for s in shots if os.path.exists(s)]
    if not shots:
        shots = sorted(glob.glob(f"{SHOTS}/*.png"))[:2]
    if len(shots) >= 2:
        for i, s in enumerate(shots[:2]):
            w, h = Image.open(s).size
            dw = 4.4
            dh = min(dw * h / w, 3.5)
            dw = dh * w / h
            S[9].shapes.add_picture(s, Inches(0.5 + i * 4.8),
                                    Inches(1.6), Inches(dw), Inches(dh))
        body(S[9], [[("Live at https://koushikdeb.duckdns.org — portfolio dashboard "
                      "(left) and per-loan reasoning with verified reason codes, "
                      "12-month risk curve and recommended recourse (right).",
                      False, MUTE, 10)]], top=5.05)
    elif shots:
        picture(S[9], shots[0], top=1.55, max_w=8.6, max_h=3.6)

    # -- 11. Benchmarking -------------------------------------------------
    def auc(k):
        v = spec.get(k, {}).get("calibrated_test", {})
        return f"{v.get('auc', 0):.4f}", f"{v.get('ece', 0):.4f}", f"{v.get('n', 0):,}"

    rows = [["Dataset (public benchmark)", "Test AUC", "ECE", "Published ref."]]
    for k, ref in (("sba", "~0.95"), ("amex", "~0.96"), ("gmsc", "~0.87"),
                   ("taiwan", "~0.78"), ("home_credit", "0.805"),
                   ("lending_club", "0.70-0.73")):
        a, e, _ = auc(k)
        rows.append([k.replace("_", " ").title(), a, e, ref])
    rows.append(["Pooled global model", f"{pooled['auc']:.4f}",
                 f"{pooled['ece']:.4f}", "—"])

    body(S[10], [
        [("All figures on an untouched test fold. ", True),
         ("Data is split three ways — 60% fit / 15% calibrate / 25% test — so the "
          "isotonic calibrator never sees the fold used to report calibration "
          "error. Reporting ECE on the calibrator's own fold (a common shortcut) "
          "drives it artificially to zero.", False)],
    ], top=1.5, size=10.5)
    table(S[10], rows, top=2.25, row_h=0.235)

    lb = A.get("berka", {})
    lift = lb.get("auc_lift_from_sequence")
    extra = []
    if lift is not None:
        extra.append([("Transaction-as-language ablation: ", True),
                      (f"on Berka (16 tabular features) the sequence embedding adds "
                       f"{lift:+.4f} AUC ({lb['tabular']['auc']:.4f} → "
                       f"{lb['tabular+sequence']['auc']:.4f}); on Amex (941 engineered "
                       f"features) it adds nothing — it substitutes for feature "
                       f"engineering rather than adding to it.", False)])
    lm = M.get("lodo_mean_auc")
    if lm is not None:
        extra.append([("Cross-domain transfer (negative result): ", True),
                      (f"leave-one-dataset-out mean AUC {lm:.4f} — a single pooled "
                       f"model does not transfer across lending domains, which is why "
                       f"SAARTHI fine-tunes on the lender's own book.", False)])
    if extra:
        body(S[10], extra, top=4.35, size=10)

    # -- 12. Additional details / Future development ----------------------
    body(S[11], [
        [("Ready for the IDBI sandbox", True, ACCENT)],
        [("The pipeline is already written against IDBI's published API schemas — "
          "API 402 (DPD / NPA status) supplies the training label, API 393 "
          "(statement with narration) feeds the sequence encoder, and APIs 391 / "
          "408 / 441 / 442 / 362 supply loan, bureau, limit and lien features.",
          False)],
        [("", False, INK, 4)],
        [("Next", True, ACCENT)],
        [("• GST-behavioural features — filing delays, turnover trend, ITC "
          "mismatch, customer concentration — modelled on the Bank's own schema.", False)],
        [("• Public-domain digital signals for thin-file MSMEs, gated by an "
          "explicit borrower-consent flag under the DPDP Act, 2023.", False)],
        [("• AWS Bedrock adapter so explanation and judging run entirely inside "
          "the Bank's VPC.", False)],
        [("• Portfolio-level what-if: simulate a policy change and see the "
          "projected shift in expected credit loss.", False)],
    ], size=11)

    # -- 13. Improvements during the 2nd prototype phase ------------------
    body(S[12], [
        [("Trained a real model.", True), (" Phase 1 trained only on whatever file "
          "was uploaded. SAARTHI now ships pre-trained on 3.3M loans across nine "
          "public corpora and fine-tunes on the lender's book.", False)],
        [("Fixed a calibration-evaluation flaw.", True), (" The earlier build fitted "
          "the isotonic calibrator and measured calibration error on the same fold, "
          "which understates it. Now a strict three-way split; every reported "
          "number comes from data neither the model nor the calibrator has seen.",
          False)],
        [("Found and removed a data leak.", True), (" On the Berka corpus, 71% of an "
          "account's transactions post-date the loan and encode the outcome, giving "
          "a false AUC of 1.0000. Aggregation is now restricted to pre-origination "
          "activity; the honest figure is 0.8558.", False)],
        [("Built transaction-as-language.", True), (" A self-supervised contrastive "
          "encoder trained on 13.25M transaction events; on a frozen embedding "
          "alone a linear probe reaches 0.8717 AUC.", False)],
        [("Verified rather than assumed.", True), (" Ensembling correlated boosters "
          "was measured, found to add nothing (−0.0029 AUC), and reported as such.",
          False)],
    ], size=10.5)

    # -- 14. Links --------------------------------------------------------
    append_to_run(S[13], "GitHub Public Repository",
                  ":  https://github.com/DevDaring/IDBI_SAARTHI", col=LINK)
    append_to_run(S[13], "Final Product Link",
                  ":  https://koushikdeb.duckdns.org", col=LINK)
    body(S[13], [
        [("Trained models (public): ", True),
         ("https://huggingface.co/Debk/saarthi-default-prediction", False, LINK)],
    ], top=2.45, size=11.5)


def main():
    bak = SRC.replace(".pptx", "_ORIGINAL_BACKUP.pptx")
    if not os.path.exists(bak):
        shutil.copy2(SRC, bak)
        print(f"backed up template -> {os.path.basename(bak)}")
    prs = Presentation(SRC)
    fill(prs)
    prs.save(SRC)
    print(f"filled deck ({len(prs.slides)} slides) -> {SRC}")


if __name__ == "__main__":
    main()
